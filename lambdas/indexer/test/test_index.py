"""
Tests for the ES indexer. This function consumes events from SQS.
"""
import datetime
import io
import json
import os
from copy import deepcopy
from gzip import compress
from io import BytesIO
from math import floor
from pathlib import Path
from string import ascii_lowercase
from time import time
from unittest import TestCase
from unittest.mock import ANY, patch
from urllib.parse import unquote_plus

import boto3
import botocore
import pptx
import pytest
import responses
from botocore import UNSIGNED
from botocore.client import Config
from botocore.stub import Stubber
from dateutil.tz import tzutc
from document_queue import EVENT_PREFIX, RetryError

from quilt_shared.const import MANIFESTS_PREFIX, NAMED_PACKAGES_PREFIX
from quilt_shared.es import (
    PACKAGE_INDEX_SUFFIX,
    get_manifest_doc_id,
    get_ptr_doc_id,
)
from t4_lambda_shared.utils import separated_env_to_iter

from .. import document_queue, index

BASE_DIR = Path(__file__).parent / 'data'

CREATE_EVENT_TYPES = {
    "ObjectCreated:Put",
    "ObjectCreated:Copy",
    "ObjectCreated:Post",
    "ObjectCreated:CompleteMultipartUpload"
}
DELETE_EVENT_TYPES = {
    "ObjectRemoved:Delete",
    "ObjectRemoved:DeleteMarkerCreated"
}
ES_ENVIRONMENT = {
    'ES_ENDPOINT': 'https://example.com',
    'AWS_ACCESS_KEY_ID': 'test_key',
    'AWS_SECRET_ACCESS_KEY': 'test_secret',
    'AWS_DEFAULT_REGION': 'ng-north-1',
}
UNKNOWN_EVENT_TYPE = "Event:WeNeverHeardOf"
# Reshaped EventBridge events from CloudTrail => S3 (see index.py notes)
EVENTBRIDGE_CORE = {
    'awsRegion': 'us-east-1',
    'eventName': 'PutObject',
    'eventTime': '2020-12-30T21:29:41Z',
    's3': {
        'bucket': {
            'name': 'quilt-s3-eventbridge'
        },
        'object': {
            'isDeleteMarker': '',
            'key': 'clear/README.md',
            'versionId': ''
        }
    }
}
# See for event structure:
# https://docs.aws.amazon.com/AmazonS3/latest/dev/notification-content-structure.html
EVENT_CORE = {
    "awsRegion": "us-east-1",
    "eventName": "ObjectCreated:Put",
    "eventSource": "aws:s3",
    "eventTime": "2020-05-22T00:32:20.515Z",
    "eventVersion": "2.1",
    "requestParameters": {"sourceIPAddress": "127.0.0.1"},
    "responseElements": {
        "x-amz-id-2": "EXAMPLE123/5678abcdefghijklambdaisawesome/mnopqrstuvwxyzABCDEFGH",
        "x-amz-request-id": "EXAMPLE123456789"
    },
    "s3": {
        "bucket": {
            "arn": "arn:aws:s3:::test-bucket",
            "name": "test-bucket",
            "ownerIdentity": {
                "principalId": "EXAMPLE"
            }
        },
        "configurationId": "testConfigRule",
        "object": {
            "key": "hello+world.txt",
            "sequencer": "0A1B2C3D4E5F678901"
        },
        "s3SchemaVersion": "1.0"
    },
    "userIdentity": {"principalId": "EXAMPLE"}
}
# typical response to head_object or get_object
OBJECT_RESPONSE = {
    'ResponseMetadata': ANY,
    'AcceptRanges': 'bytes',
    'LastModified': datetime.datetime(2019, 5, 30, 23, 27, 29, tzinfo=tzutc()),
    'ContentLength': 1555,
    'ETag': '"8dbf7b98d5458a46327fb58f27b9af6e"',
    'VersionId': 'wcOZpjy5G.tJ2N.rwPhiR.NY_RftJ3A_',
    'ContentType': 'binary/octet-stream',
    'Metadata': {
        'helium': '{"user_meta": {}}'
    }
}
# Simulated first line of manifest
MANIFEST_DATA = {
    "version": "v0",
    "user_meta": {
        "arbitrary": "metadata",
        "list": [5, 9, 19],
        "int": 42,
        "object": {"treble": "is", "a": "friend"}
    },
    "message": "interesting comment with interesting symbols #$%@☮ 😎!"
}


def _check_event(synthetic, organic):
    # Ensure that synthetic events have the same shape as actual organic ones,
    # and that overridden properties like bucket, key, eTag are properly set
    # same keys at top level
    assert organic.keys() == synthetic.keys()
    # same value types (values might differ and that's OK)
    assert {type(v) for v in organic.values()} == \
        {type(v) for v in synthetic.values()}
    # same keys and nested under "s3"
    assert organic["s3"].keys() == synthetic["s3"].keys()
    assert organic["s3"]["bucket"].keys() == synthetic["s3"]["bucket"].keys()
    # same value types under S3 (values might differ and that's OK)
    assert {type(v) for v in organic["s3"].values()} == \
        {type(v) for v in synthetic["s3"].values()}
    # spot checks for overridden properties
    # size absent on delete
    if "size" in organic["s3"]["bucket"]:
        assert organic["s3"]["object"]["size"] == synthetic["s3"]["object"]["size"]
    # versionId absent when unversioned bucket or hard delete
    if "versionId" in organic["s3"]["bucket"]:
        assert organic["s3"]["object"]["versionId"] == synthetic["s3"]["object"]["versionId"]
    # should always be present
    assert organic["awsRegion"] == synthetic["awsRegion"]
    assert organic["s3"]["bucket"]["arn"] == synthetic["s3"]["bucket"]["arn"]
    assert organic["s3"]["bucket"]["name"] == synthetic["s3"]["bucket"]["name"]
    assert organic["s3"]["object"]["key"] == synthetic["s3"]["object"]["key"]
    assert organic["s3"]["object"]["eTag"] == synthetic["s3"]["object"]["eTag"]


def make_event(
        name,
        *,
        bucket="test-bucket",
        eTag="123456",
        key="hello+world.txt",
        region="us-east-1",
        size=100,
        versionId="1313131313131.Vier50HdNbi7ZirO65",
        bucket_versioning=True
):
    """create an event based on EVENT_CORE, add fields to match organic AWS events"""
    if name in CREATE_EVENT_TYPES:
        args = {
            "bucket": bucket,
            "eTag": eTag,
            "key": key,
            "region": region,
            "size": size
        }
        if bucket_versioning:
            args["versionId"] = versionId
        return _make_event(
            name,
            **args
        )
    # no versionId or eTag in this case
    elif name == "ObjectRemoved:Delete":
        return _make_event(
            name,
            bucket=bucket,
            key=key,
            region=region
        )
    elif name == "ObjectRemoved:DeleteMarkerCreated":
        # these events are possible in both versioned and unversioned buckets
        # (e.g. bucket now unversioned that was versioned will generate a
        # delete marker on `aws s3 rm`)
        args = {
            "bucket": bucket,
            "eTag": eTag,
            "key": key,
            "region": region,
            "size": size
        }
        if bucket_versioning:
            args["versionId"] = versionId
        return _make_event(
            name,
            **args
        )
    elif name == UNKNOWN_EVENT_TYPE:
        return _make_event(UNKNOWN_EVENT_TYPE)

    else:
        raise ValueError(f"Unexpected event type: {name}")


def _make_event(
        name,
        *,
        bucket="",
        eTag="",
        key="",
        region="",
        size=0,
        versionId=""
):
    """make events in the pattern of
    https://docs.aws.amazon.com/AmazonS3/latest/dev/notification-content-structure.html
    and
    AWS Lambda > Console > Test Event
    """
    e = deepcopy(EVENT_CORE)
    e["eventName"] = name

    if bucket:
        e["s3"]["bucket"]["name"] = bucket
        e["s3"]["bucket"]["arn"] = f"arn:aws:s3:::{bucket}"
    if key:
        e["s3"]["object"]["key"] = key
    if eTag:
        e["s3"]["object"]["eTag"] = eTag
    if size:
        e["s3"]["object"]["size"] = size
    if region:
        e["awsRegion"] = region
    if versionId:
        e["s3"]["object"]["versionId"] = versionId

    return e


@pytest.mark.parametrize(
    "event_type, kwargs",
    [
        (
            "ObjectRemoved:Delete",
            {
                "bucket": "test",
                "etag": "123",
                "ext": ".txt",
                "key": "foo",
                "last_modified": datetime.datetime(2019, 5, 30, 23, 27, 29, tzinfo=tzutc()),
                "version_id": "abc",
                "size": 0,
                "text": "",
                "s3_tags": None,
            }
        ),
        (
            "ObjectCreated:Copy",
            {
                "bucket": "test",
                "etag": "123",
                "ext": ".jsonl",
                "key": "foo",
                "last_modified": datetime.datetime(2019, 5, 30, 23, 27, 29, tzinfo=tzutc()),
                "size": 0,
                "text": "iajsoeqroieurqwiuro•",
                "version_id": "abc",
                "s3_tags": {"key": "value"},
            }
        ),
    ]
)
@patch.object(index.DocumentQueue, 'append_document')
def test_append(_append_mock, event_type, kwargs):
    """test document_queue.append; outside of class so we can parameterize"""
    dq = index.DocumentQueue(None)
    dq.append(event_type=event_type, **kwargs)
    assert _append_mock.call_count == 1


@pytest.mark.parametrize(
    "ext, expected", [
        (".gz", "gz"),
        (".bz2", None),
        (".zip", None),
        (".csv", None),
    ]
)
def test_get_compression(ext, expected):
    assert index.get_compression(ext) == expected


@pytest.mark.parametrize(
    "key, expected", (
        ("foo/bar.baz/hello world.txt", ("", ".txt")),
        ("foo/bar.baz/hello.world.csv", (".world", ".csv")),
        ("foo/bar.baz/hello-world.csv.gz", (".csv", ".gz")),
        ("abarefilenoextensions", ("", "")),
        ("f", ("", "")),
        (" ", ("", "")),
        ("..", ("", "")),
        (".", ("", "")),
        ("", ("", "")),
    )
)
def test_get_normalized_extensions(key, expected):
    assert index.get_normalized_extensions(key) == expected


@pytest.mark.parametrize(
    "key, compression, expected", [
        # parquet
        ("s3/some/file.c000", None, ".parquet"),
        # parquet, nonzero part number
        ("s3/some/file.c001", None, ".parquet"),
        # -c0001 file
        ("s3/some/file-c0001", None, ".parquet"),
        # -c00111 file (should never happen)
        ("s3/some/file-c000121", None, ""),
        ("s3/some/file-boom-boom!bam.pq", None, ".parquet"),
        ("s3/some/file-boom-boom!bam.pq.gz", "gz", ".parquet"),
        pytest.param("s3/some/file-boom-boom!bam.pq.gz", None, ".parquet", marks=pytest.mark.xfail),
        ("s3/some/file-boom-boom!maga_0", None, ".parquet"),
        # .txt file, should be unchanged
        ("s3/some/file-c0000.txt", None, ".txt"),
        ("s3/to.path/a/file.with.lots.of.dots.h5ad.json", None, ".json"),
        ("s3/to.path/a/file.csv.gz", "gz", ".csv"),
        # we dont' support .bz2
        ("s3/to.path/a/file.csv.bz2", None, ".bz2")
    ]
)
def test_infer_extensions(key, compression, expected):
    """ensure we are guessing file types well"""
    assert index.infer_extensions(key, index.get_normalized_extensions(key), compression) == expected


def test_map_event_name_and_validate():
    """ensure that we map eventName properly, ensure that shape validation code works"""
    for name in CREATE_EVENT_TYPES.union(DELETE_EVENT_TYPES).union({UNKNOWN_EVENT_TYPE}):
        event = make_event(name)
        assert index.shape_event(event)
        original = event["eventName"]
        assert original == index.map_event_name(event)

    for name in index.EVENTBRIDGE_TO_S3:
        event = EVENTBRIDGE_CORE.copy()
        event["eventName"] = name
        mapped = index.map_event_name(event)
        assert name != mapped
        if name == "DeleteObject":
            assert mapped.startswith(EVENT_PREFIX["Removed"])
        else:
            assert mapped.startswith(EVENT_PREFIX["Created"])

        assert index.shape_event(event)

    delete_marker = EVENTBRIDGE_CORE.copy()
    delete_marker["eventName"] = "DeleteObject"
    delete_marker["s3"]["object"]["isDeleteMarker"] = "true"
    delete_marker["s3"]["object"]["versionId"] = "someid"
    mapped = index.map_event_name(delete_marker)
    assert mapped == "ObjectRemoved:DeleteMarkerCreated"
    reshaped = index.shape_event(delete_marker.copy())
    assert reshaped
    assert delete_marker != reshaped

    assert not index.shape_event({"bad": "event"})

    malformed = EVENTBRIDGE_CORE.copy()
    del malformed["s3"]["bucket"]["name"]
    assert not index.shape_event(malformed)


@pytest.mark.parametrize(
    "env_var, check, expected", [
        (".txt,.csv", ".parquet", False),
        (".txt,.csv", ".csv", True),
        (".txt,.csv", ".txt", True),
        (".parquet,.tsvl", ".parquet", True),
        (".parquet,.tsvl", ".csv", False),
    ]
)
def test_skip_rows_env(env_var, check, expected):
    """test whether or not index skips rows per SKIP_ROWS_EXTS=LIST"""
    # because of module caching we can't just patch the environment variable
    # since index.SKIP_ROWS_EXTS will never change after import
    with patch.dict(os.environ, {'SKIP_ROWS_EXTS': env_var}):
        exts = separated_env_to_iter('SKIP_ROWS_EXTS')
        with patch('index.SKIP_ROWS_EXTS', exts):
            if expected:
                assert check in exts
            else:
                assert check not in exts


class MockContext():
    def get_remaining_time_in_millis(self):
        return 30000


class TestIndex(TestCase):
    def setUp(self):
        # total number of times we expect the ES _bulk API is called
        # we do not use `len(responses.calls)` because it always evaluates to 0
        # during both setup and teardown; reason is that we are using add_callback()?
        self.actual_es_calls = 0
        self.requests_mock = responses.RequestsMock(assert_all_requests_are_fired=True)
        self.requests_mock.start()

        # Create a dummy S3 client that (hopefully) can't do anything.
        self.s3_client = boto3.client('s3', config=Config(signature_version=UNSIGNED))

        self.s3_client_patcher = patch(
            __name__ + '.index.make_s3_client',
            return_value=self.s3_client
        )
        self.s3_client_patcher.start()

        self.s3_stubber = Stubber(self.s3_client)
        self.s3_stubber.activate()

        self.env_patcher = patch.dict(os.environ, ES_ENVIRONMENT)
        self.env_patcher.start()

    def tearDown(self):
        self.env_patcher.stop()

        self.s3_stubber.assert_no_pending_responses()
        self.s3_stubber.deactivate()
        self.s3_client_patcher.stop()

        self.requests_mock.stop()

    def _get_contents(self, name, ext, compression=None, size=123):
        return index.maybe_get_contents(
            'test-bucket',
            name,
            ext,
            etag='etag',
            version_id=None,
            s3_client=self.s3_client,
            size=size,
            compression=compression
        )

    def _make_es_callback(
            self,
            *,
            errors=False,
            status=200,
            unknown_items=False
    ):
        """
        create a callback that checks the shape of the response
        TODO: handle errors and delete actions
        """
        def check_response(request):
            raw = [json.loads(line) for line in request.body.splitlines()]
            # drop the optional source and isolate the actions
            # see https://www.elastic.co/guide/en/elasticsearch/reference/6.7/docs-bulk.html
            actions = [line for line in raw if len(line.keys()) == 1]
            items = [
                {
                    top_key: {
                        "_id": values["_id"],
                        "_index": values["_index"],
                        "_type": "_doc",
                        "status": 200
                    }
                }
                for action in actions
                for top_key, values in action.items()
            ]
            if unknown_items:
                items = [
                    {"event_we_never_heard_of": value}
                    for item in items
                    for value in item.values()
                ]
            # see https://www.elastic.co/guide/en/elasticsearch/reference/6.7/docs-bulk.html
            # for response format
            response = {
                "took": 5*len(actions),
                "errors": errors,
                "items": items
            }
            self.actual_es_calls = self.actual_es_calls + 1

            return (status, {}, json.dumps(response))

        return check_response

    def _test_index_events(
            self,
            event_names,
            *,
            bucket_versioning=True,
            errors=False,
            expected_es_calls=0,
            mock_elastic=True,
            mock_overrides=None,
            status=200,
            unknown_items=False,

    ):
        """
        Reusable helper function to test indexing files based on on or more
        events
        """
        inner_records = []
        for name in event_names:
            event_kwargs = mock_overrides.get('event_kwargs', {}) if mock_overrides else {}
            event = make_event(name, bucket_versioning=bucket_versioning, **event_kwargs)
            inner_records.append(event)
            now = index.now_like_boto3()
            un_key = unquote_plus(event["s3"]["object"]["key"])
            eTag = event["s3"]["object"].get("eTag")
            versionId = event["s3"]["object"].get("versionId")

            expected_params = {
                'Bucket': event["s3"]["bucket"]["name"],
                'Key': un_key,
            }
            # We only get versionId for certain events (when bucket versioning is
            # on (or was on and a delete-object is issued with a particular version-id?)
            if versionId:
                expected_params["VersionId"] = versionId
            elif eTag:
                expected_params["IfMatch"] = eTag
            # infer mock status (we only talk head S3 on create events)
            mock_get_object_tagging = mock_head = mock_object = name in CREATE_EVENT_TYPES
            # check for occasional overrides (which can be false)
            if mock_overrides and "mock_head" in mock_overrides:
                mock_head = mock_overrides.get("mock_head")
            if mock_overrides and "mock_object" in mock_overrides:
                mock_object = mock_overrides.get("mock_object")

            if mock_head:
                self.s3_stubber.add_response(
                    method='head_object',
                    service_response={
                        'Metadata': {},
                        'ContentLength': event["s3"]["object"]["size"],
                        'LastModified': now,
                    },
                    expected_params=expected_params
                )

            if mock_object:
                if mock_overrides and mock_overrides.get('skip_byte_range'):
                    expected = expected_params.copy()
                else:
                    expected = {
                        **expected_params,
                        'Range': 'bytes=0-99'
                    }
                self.s3_stubber.add_response(
                    method='get_object',
                    service_response={
                        'Metadata': {},
                        'ContentLength': event["s3"]["object"]["size"],
                        'LastModified': now,
                        'Body': BytesIO(b'Hello World!'),
                    },
                    expected_params=expected
                )

            if mock_get_object_tagging:
                expected = {
                    "Bucket": event["s3"]["bucket"]["name"],
                    "Key": un_key,
                }
                if versionId:
                    expected["VersionId"] = versionId
                self.s3_stubber.add_response(
                    method="get_object_tagging",
                    service_response={
                        "TagSet": [
                            {"Key": "key", "Value": "value"},
                        ]
                    },
                    expected_params=expected,
                )

        if mock_elastic:
            self.requests_mock.add_callback(
                responses.POST,
                'https://example.com:443/_bulk',
                callback=self._make_es_callback(
                    errors=errors,
                    status=status,
                    unknown_items=unknown_items
                ),
                content_type='application/json'
            )

        records = {
            "Records": [{
                "body": json.dumps({
                    "Message": json.dumps({
                        "Records": inner_records
                    })
                })
            }]
        }

        index.handler(records, MockContext())
        assert self.actual_es_calls == expected_es_calls, \
            (
                f"Expected ES endpoint to be called {expected_es_calls} times, "
                f"got {self.actual_es_calls} calls instead"
            )

    @patch.object(index.DocumentQueue, 'send_all')
    @patch.object(index.DocumentQueue, 'append')
    @patch.object(index, 'maybe_get_contents')
    @patch.object(index, 'index_if_pointer')
    def test_40X(self, index_if_mock, contents_mock, append_mock, send_mock):
        """
        test fatal head 40Xs that will cause us to skip objects
        """
        bucket = "somebucket"
        key = "events/copy-one/0.ext"
        etag = "7b4b71116bb21d3ea7138dfe7aabf036"
        version_ids = ["Yj1vyLWcE9FTFIIrsgk.yAX7NbJrAW7g", "null"]

        error_codes = ["402", "403", "404"]

        for version_id in version_ids:
            for error_code in error_codes:
                self.s3_stubber.add_client_error(
                    method="head_object",
                    service_error_code=error_code,
                    service_message=f"An error occurred ({error_code}) when calling the HeadObject operation",
                    http_status_code=int(error_code),
                    expected_params={
                        "Bucket": bucket,
                        "Key": key,
                        "VersionId": version_id,
                    },
                )
                # in the 403 case we try again without VersionId
                if error_code == "403" and version_id == "null":
                    self.s3_stubber.add_client_error(
                        method="head_object",
                        service_error_code=error_code,
                        service_message=f"An error occurred ({error_code}) when calling the HeadObject operation",
                        http_status_code=int(error_code),
                        expected_params={
                            "Bucket": bucket,
                            "Key": key,
                            "IfMatch": etag,
                        },
                    )
                event = make_event(
                    "ObjectCreated:Put",
                    bucket=bucket,
                    key=key,
                    size=73499,
                    eTag=etag,
                    region="us-west-1",
                    versionId=version_id,
                )

                records = {
                    "Records": [{
                        "body": json.dumps({
                            "Message": json.dumps({
                                "Records": [event]
                            })
                        })
                    }]
                }

                index.handler(records, None)

                assert index_if_mock.call_count == 0
                assert contents_mock.call_count == 0
                assert append_mock.call_count == 0

        assert send_mock.call_count == len(error_codes)*len(version_ids)

    def test_create_event_failure(self):
        """
        Check that the indexer doesn't blow up on create event failures.
        """
        with pytest.raises(RetryError, match="Failed to load"):
            self._test_index_events(
                ["ObjectCreated:Put"],
                errors=True,
                status=400
            )

    def test_create_copy_index(self):
        """test indexing a single file from copy event"""
        self._test_index_events(
            ["ObjectCreated:Copy"],
            expected_es_calls=1
        )
        # Elastic only needs to be mocked once per test

    def test_create_put_index(self):
        """test indexing a single file from put event"""
        self._test_index_events(
            ["ObjectCreated:Put"],
            expected_es_calls=1
        )

    def test_create_put_index_unversioned(self):
        """test indexing a single file from put event"""
        self._test_index_events(
            ["ObjectCreated:Put"],
            bucket_versioning=False,
            expected_es_calls=1
        )

    def test_create_post_index(self):
        """test indexing a single file from post event"""
        self._test_index_events(
            ["ObjectCreated:Post"],
            expected_es_calls=1
        )

    def test_create_multipart_index(self):
        """test indexing a single file from post event"""
        self._test_index_events(
            ["ObjectCreated:CompleteMultipartUpload"],
            expected_es_calls=1
        )

    def test_delete_event(self):
        """
        Check that the indexer doesn't blow up on delete events.
        """
        self._test_index_events(
            ["ObjectRemoved:Delete"],
            expected_es_calls=1
        )

    def test_delete_event_failure(self):
        """
        Check that the indexer doesn't blow up on delete event failures.
        """
        # TODO, why does pytest.raises(RetryError not work?)
        with pytest.raises(RetryError, match="Failed to load"):
            self._test_index_events(
                ["ObjectRemoved:Delete"],
                errors=True,
                status=400
            )

    def test_delete_event_no_versioning(self):
        """
        Check that the indexer doesn't blow up on delete events.
        """
        self._test_index_events(
            ["ObjectRemoved:Delete"],
            bucket_versioning=False,
            expected_es_calls=1
        )

    def test_delete_marker_event(self):
        """
        common event in versioned buckets
        """
        self._test_index_events(
            ["ObjectRemoved:DeleteMarkerCreated"],
            expected_es_calls=1
        )

    def test_delete_marker_event_no_versioning(self):
        """
        this can happen if a bucket was versioned, and now isn't, followed by
        `aws s3 rm`
        """
        # don't mock head or get; this event should never call them
        self._test_index_events(
            ["ObjectRemoved:DeleteMarkerCreated"],
            bucket_versioning=False,
            expected_es_calls=1
        )

    def test_extract_pdf(self):
        """test pdf extraction to text"""
        with open(BASE_DIR / "MUMmer.pdf", "rb") as pdf:
            txt = index.extract_pdf(pdf)
            phrases = [
                "Alignment of whole genomes",
                "When the genome sequence of two closely related organisms",
                # 2 lines as one string
                "; the result is a very detailed and inclusive base-to-base mapping "
                "between the two sequences.",
                # 4 lines as one string
                "Although our alignment does not contain all the details generated "
                "and displayed by the combination of methods used in Ansari-Lari "
                "et al., the overall alignment of the two sequences is easily "
                "apparent from the output of our program.",
                "under Grant no. R01-AI40125-01.",
            ]
            assert all(p in txt for p in phrases)

    @patch.object(index, 'extract_parquet')
    def test_index_c000(self, extract_mock):
        """ensure files with special extensions get treated as parquet"""
        extract_mock.return_value = ('parquet-body', {'schema': {'names': []}})
        self._test_index_events(
            ["ObjectCreated:Put"],
            expected_es_calls=1,
            mock_overrides={
                "event_kwargs": {
                    "key": "obscure_path/long-complicated-name-c000",
                },
                # no byte ranges for parquet files
                "skip_byte_range": True
            }
        )
        extract_mock.assert_called_once()

    @patch.object(index.DocumentQueue, 'append')
    @patch.object(index, 'maybe_get_contents')
    def test_index_c000_contents(self, get_mock, append_mock):
        """ensure files with special extensions get treated as parquet"""
        parquet_data = b'@@parquet-data@@'
        get_mock.return_value = parquet_data
        self._test_index_events(
            ["ObjectCreated:Put"],
            # we're mocking append so ES will never get called
            mock_elastic=False,
            mock_overrides={
                "event_kwargs": {
                    # this key should infer to parquet
                    "key": "obscure_path/long-complicated-name-c000"
                },
                # we patch get_contents so _test_index_events doesn't need to
                "mock_object": False,
                # no byte ranges for parquet files
                "skip_byte_range": True
            }
        )
        get_mock.assert_called_once()
        # ensure parquet data is getting to elastic
        append_mock.assert_called_once_with(
            event_type='ObjectCreated:Put',
            bucket='test-bucket',
            etag='123456',
            ext='',
            key='obscure_path/long-complicated-name-c000',
            last_modified=ANY,
            size=100,
            text=parquet_data,
            version_id='1313131313131.Vier50HdNbi7ZirO65',
            s3_tags={"key": "value"},
        )

    @patch.object(index, 'maybe_get_contents')
    def test_index_exception(self, get_mock):
        """test indexing a single file that throws an exception"""
        class ContentException(Exception):
            pass
        get_mock.side_effect = ContentException("Unable to get contents in Glacier")
        # get_mock already mocks get_object, so don't mock it in _test_index_event
        self._test_index_events(
            ["ObjectCreated:Put"],
            expected_es_calls=1,
            mock_overrides={
                "mock_object": False
            }
        )

    @patch.object(index, "es")
    @patch.object(index.DocumentQueue, "append_document")
    def test_index_if_pointer_not_exists(self, append_mock, es_mock):
        bucket = "quilt-example"
        key = f"{NAMED_PACKAGES_PREFIX}author/semantic/1610412903"

        self.s3_stubber.add_client_error(
            method="get_object",
            expected_params={
                "Bucket": bucket,
                "Key": key,
            },
            service_error_code="NoSuchKey",
        )

        index.index_if_pointer(
            self.s3_client,
            index.DocumentQueue(None),
            bucket=bucket,
            key=key,
        )

        es_mock.delete_by_query.assert_called_once_with(
            index=bucket + PACKAGE_INDEX_SUFFIX,
            body={
                "query": {
                    "bool": {
                        "filter": [{"term": {"_id": get_ptr_doc_id("author/semantic", "1610412903")}}],
                    }
                }
            },
        )
        append_mock.assert_not_called()

    @patch.object(index, "es")
    @patch.object(index.DocumentQueue, 'append_document')
    def test_index_if_pointer(self, append_mock, es_mock):
        bucket = "quilt-example"
        handle = "author/semantic"
        pointer_file = "1610412903"
        key = f"{NAMED_PACKAGES_PREFIX}{handle}/{pointer_file}"
        pkg_hash = "a" * 64
        last_modified = datetime.datetime(2021, 1, 1, 0, 0, tzinfo=tzutc())

        self.s3_stubber.add_response(
            method="get_object",
            expected_params={
                "Bucket": bucket,
                "Key": key,
            },
            service_response={
                "Body": BytesIO(pkg_hash.encode()),
                "LastModified": last_modified,
            },
        )

        index.index_if_pointer(
            self.s3_client,
            index.DocumentQueue(None),
            bucket=bucket,
            key=key,
        )

        es_mock.delete_by_query.assert_called_once_with(
            index=bucket + PACKAGE_INDEX_SUFFIX,
            body={
                "query": {
                    "bool": {
                        "filter": [{"term": {"_id": get_ptr_doc_id("author/semantic", "1610412903")}}],
                        "must_not": [{"term": {"join_field#mnfst": get_manifest_doc_id("a" * 64)}}],
                    }
                }
            },
        )
        append_mock.assert_called_once_with({
            "_index": bucket + PACKAGE_INDEX_SUFFIX,
            "_op_type": "index",
            "_id": get_ptr_doc_id(handle, pointer_file),
            "join_field": {
                "name": "ptr",
                "parent": get_manifest_doc_id(pkg_hash),
            },
            "routing": get_manifest_doc_id(pkg_hash),
            "ptr_name": handle,
            "ptr_tag": pointer_file,
            "ptr_last_modified": last_modified,
        })

    def test_index_if_pointer_skip(self):
        """test cases where index_if_pointer ignores input for different reasons"""
        # none of these should index due to out-of-range timestamp or non-integer name
        for file_name in [1451631500, 1767250801]:
            key = f".quilt/named_packages/foo/bar/{file_name}"
            assert not index.index_if_pointer(
                self.s3_client,
                index.DocumentQueue(None),
                bucket="quilt-example",
                key=key,
            )
        # none of these should index due to bad file path
        good_timestamp = floor(time())
        for key in [
                f".quilt/named_packages//{good_timestamp}.txt",
                f".quilt/named_packages/{good_timestamp}",
                f".quilt/named_packages/not-deep-enough/{good_timestamp}",
                ".quilt/named_packages/usr/pkg/",
                ".quilt/named_packages/usr/",
                ".quilt/named_packages/",
                f"somewhere/else/foo/bar/{floor(time())}",
        ]:
            assert not index.index_if_pointer(
                self.s3_client,
                index.DocumentQueue(None),
                bucket="quilt-example",
                key=key,
            )

    @patch.object(index.DocumentQueue, 'append')
    @patch.object(index, 'maybe_get_contents')
    @patch.object(index, 'index_if_pointer')
    def test_index_if_pointer_negative(self, index_mock, get_mock, append_mock):
        """test non-manifest file (still calls index_if_pointer)"""
        json_data = json.dumps({"version": 1})
        get_mock.return_value = json_data

        self._test_index_events(
            ["ObjectCreated:Put"],
            # we're mocking append so ES will never get called
            mock_elastic=False,
            mock_overrides={
                "event_kwargs": {
                    "key": "obscure_path/long-complicated-name-c000"
                },
                # we patch maybe_get_contents so _test_index_events doesn't need to
                "mock_object": False,
                # no byte ranges for parquet files
                "skip_byte_range": True
            }
        )
        get_mock.assert_called_once()
        index_mock.assert_called_once()
        append_mock.assert_called_once_with(
            event_type='ObjectCreated:Put',
            bucket='test-bucket',
            etag='123456',
            ext='',
            key='obscure_path/long-complicated-name-c000',
            last_modified=ANY,
            size=100,
            text=json_data,
            version_id='1313131313131.Vier50HdNbi7ZirO65',
            s3_tags={"key": "value"},
        )

    @patch.object(index.sqs, "send_message")
    def test_manifest_event(self, send_message_mock):
        """test indexing a manifest file"""
        self._test_index_events(
            ["ObjectCreated:Put"],
            expected_es_calls=1,
            mock_overrides={
                "event_kwargs": {
                    "key": MANIFESTS_PREFIX + "a" * 64,
                },
                "mock_object": False,
            }
        )

        send_message_mock.assert_called_once()

    def test_multiple_index_events(self):
        """
        Messages from SQS contain up to N messages.
        Currently N=10, this number is determined on the backend by CloudFormation
        """
        self._test_index_events(
            [
                "ObjectCreated:Put",
                "ObjectCreated:Put",
                "ObjectCreated:Put",
                "ObjectCreated:Put",
                "ObjectCreated:Put",
                "ObjectCreated:Copy",
                "ObjectCreated:Copy",
                "ObjectCreated:Copy",
                "ObjectCreated:Copy",
                "ObjectRemoved:Delete"
            ],
            expected_es_calls=1
        )

    def test_extension_overrides(self):
        """ensure that only the file extensions in override are indexed"""
        with patch(__name__ + '.index.get_content_index_extensions', return_value={'.unique1', '.unique2'}):
            self.s3_stubber.add_response(
                method='get_object',
                service_response={
                    'Metadata': {},
                    'ContentLength': 123,
                    'Body': BytesIO(b'Hello World!'),
                },
                expected_params={
                    'Bucket': 'test-bucket',
                    'Key': 'foo.unique1',
                    'IfMatch': 'etag',
                    'Range': 'bytes=0-122',
                }
            )
            self.s3_stubber.add_response(
                method='get_object',
                service_response={
                    'Metadata': {},
                    'ContentLength': 123,
                    'Body': BytesIO(b'Hello World!'),
                },
                expected_params={
                    'Bucket': 'test-bucket',
                    'Key': 'foo.unique2',
                    'IfMatch': 'etag',
                    'Range': 'bytes=0-122',
                }
            )
            # only these two file types should be indexed
            assert self._get_contents('foo.unique1', '.unique1') == "Hello World!"
            assert self._get_contents('foo.unique2', '.unique2') == "Hello World!"
            # these files should not get content indexed, therefore no S3 mock
            assert self._get_contents('foo.txt', '.txt') == ""
            assert self._get_contents('foo.ipynb', '.ipynb') == ""

    def test_synthetic_copy_event(self):
        """check synthetic ObjectCreated:Copy event vs organic obtained on 26-May-2020
        (bucket versioning on)
        """
        synthetic = make_event(
            "ObjectCreated:Copy",
            bucket="somebucket",
            key="events/copy-one/0.png",
            size=73499,
            eTag="7b4b71116bb21d3ea7138dfe7aabf036",
            region="us-west-1",
            versionId="Yj1vyLWcE9FTFIIrsgk.yAX7NbJrAW7g"
        )
        # actual event from S3 with a few obfuscations to protect the innocent
        organic = {
            "eventVersion": "2.1",
            "eventSource": "aws:s3",
            "awsRegion": "us-west-1",
            "eventTime": "2020-05-26T22:15:10.906Z",
            "eventName": "ObjectCreated:Copy",
            "userIdentity": {
                "principalId": "AWS:EXAMPLEDUDE"
            },
            "requestParameters": {
                "sourceIPAddress": "12.999.99.999"
            },
            "responseElements": {
                "x-amz-request-id": "CEF0E4FD6D0944D7",
                "x-amz-id-2": "EXAMPLE/+GUID/m/HAApWP+3arsz0QPph7OBVdl1"
            },
            "s3": {
                "s3SchemaVersion": "1.0",
                "configurationId": "YmJkYWUyYmYtNzg5OC00NGRiLTk0NmItNDMxNzA4NzhiZDNk",
                "bucket": {
                    "name": "somebucket",
                    "ownerIdentity": {
                        "principalId": "SAMPLE"
                    },
                    "arn": "arn:aws:s3:::somebucket"
                },
                "object": {
                    "key": "events/copy-one/0.png",
                    "size": 73499,
                    "eTag": "7b4b71116bb21d3ea7138dfe7aabf036",
                    "versionId": "Yj1vyLWcE9FTFIIrsgk.yAX7NbJrAW7g",
                    "sequencer": "005ECD94EFA9B09DD8"
                }
            }
        }
        _check_event(synthetic, organic)

    def test_synthetic_copy_event_no_versioning(self):
        """check synthetic ObjectCreated:Copy event vs organic obtained on 26-May-2020
        (this is for a bucket with versioning turned off)
        """
        synthetic = make_event(
            "ObjectCreated:Copy",
            bucket="somebucket",
            key="events/copy-one-noversioning/0.png",
            size=73499,
            eTag="7b4b71116bb21d3ea7138dfe7aabf036",
            region="us-west-1",
        )
        # actual event from S3 with a few obfuscations to protect the innocent
        organic = {
            "eventVersion": "2.1",
            "eventSource": "aws:s3",
            "awsRegion": "us-west-1",
            "eventTime": "2020-05-27T20:31:45.823Z",
            "eventName": "ObjectCreated:Copy",
            "userIdentity": {"principalId": "AWS:boombomakasdfsdf"},
            "requestParameters": {"sourceIPAddress": "07.123.45.899"},
            "responseElements": {
                "x-amz-request-id": "DECF307B5F55C78D",
                "x-amz-id-2": "guid/hash/tG++guid/stuff"
            },
            "s3": {
                "s3SchemaVersion": "1.0",
                "configurationId": "stuff",
                "bucket": {
                    "name": "somebucket",
                    "ownerIdentity": {"principalId": "B3ASKDFASDFAF"},
                    "arn": "arn:aws:s3:::somebucket"},
                "object": {
                    "key": "events/copy-one-noversioning/0.png",
                    "size": 73499,
                    "eTag": "7b4b71116bb21d3ea7138dfe7aabf036",
                    "sequencer": "005ECECE336C7A4715"
                }
            }
        }
        _check_event(synthetic, organic)

    def test_synthetic_delete_marker_event_no_versioning(self):
        """check synthetic ObjectRemoved:DeleteMarkerCreated event vs organic
        obtained on 27-May-2020
        (this is for a bucket that had versioning on but now off)
        note: this is a standard delete marker minus the versionId
        """
        synthetic = make_event(
            "ObjectRemoved:DeleteMarkerCreated",
            key="events/copy-many-noversioning/0.png",
            eTag="d41d8cd98f00b204e9800998ecf8427e",
            region="us-west-1",
        )
        # actual event from S3 with a few obfuscations to protect the innocent
        organic = {
            "eventVersion": "2.1",
            "eventSource": "aws:s3",
            "awsRegion": "us-west-1",
            "eventTime": "2020-05-28T23:53:24.662Z",
            "eventName": "ObjectRemoved:DeleteMarkerCreated",
            "userIdentity": {"principalId": "AWS:boommasdfagnag"},
            "requestParameters": {"sourceIPAddress": "12.888.91.910"},
            "responseElements": {
                "x-amz-request-id": "35781DEB9DA7612E",
                "x-amz-id-2": "Qguid+Oguid+WRa/guid/guid+AwtLbBepO7QEBNbwguid/LfQguid"
            },
            "s3": {
                "s3SchemaVersion": "1.0",
                "configurationId": "guiadskfjasdlfkjasdklfjasdfd",
                "bucket": {
                    "name": "test-bucket",
                    "ownerIdentity": {"principalId": "adflkjasdklfjadf"},
                    "arn": "arn:aws:s3:::test-bucket"
                },
                "object": {
                    "key": "events/copy-many-noversioning/0.png",
                    "eTag": "d41d8cd98f00b204e9800998ecf8427e",
                    "sequencer": "005ED04EF537DAB0EE"
                }
            }
        }
        _check_event(synthetic, organic)

    def test_synthetic_put_event(self):
        """check synthetic ObjectCreated:Put event vs organic obtained on 27-May-2020
        (bucket versioning on)"""
        synthetic = make_event(
            "ObjectCreated:Copy",
            bucket="anybucket",
            key="events/put-one/storms.parquet",
            size=923078,
            eTag="502f21cfc143fb0c35f563eda5699fa9",
            region="us-west-1",
            versionId="yYSoQSg3.BfosdUxnRSv9vFg.WAPMmfn"
        )
        # actual event from S3 with a few obfuscations to protect the innocent
        organic = {
            "eventVersion": "2.1",
            "eventSource": "aws:s3",
            "awsRegion": "us-west-1",
            "eventTime": "2020-05-27T18:57:36.268Z",
            "eventName": "ObjectCreated:Put",
            "userIdentity": {"principalId": "AWS:notgonnabehereanyway"},
            "requestParameters": {"sourceIPAddress": "12.345.67.890"},
            "responseElements": {
                "x-amz-request-id": "371A83BCE4341D7D",
                "x-amz-id-2": "a+example+morestuff+01343413434234234234"
            },
            "s3": {
                "s3SchemaVersion": "1.0",
                "configurationId": "YmJkYWUyYmYtNzg5OC00NGRiLTk0NmItNDMxNzA4NzhiZDNk",
                "bucket": {
                    "name": "anybucket",
                    "ownerIdentity": {"principalId": "myidhere"},
                    "arn": "arn:aws:s3:::anybucket"
                },
                "object": {
                    "key": "events/put-one/storms.parquet",
                    "size": 923078,
                    "eTag": "502f21cfc143fb0c35f563eda5699fa9",
                    "versionId": "yYSoQSg3.BfosdUxnRSv9vFg.WAPMmfn",
                    "sequencer": "005ECEB81C34962CFC"
                }
            }
        }
        _check_event(synthetic, organic)

    def test_synthetic_multipart_event(self):
        """check synthetic ObjectCreated:Put event vs organic obtained on 27-May-2020
        (bucket versioning on)"""
        make_event(
            "ObjectCreated:CompleteMultipartUpload",
            bucket="anybucket",
            key="events/multipart-one/part-00006-495c48e6-96d6-4650-aa65-3c36a3516ddd.c000.snappy.parquet",
            size=135397292,
            eTag="0eb149127d0277326dedcf0c530ca966-17",
            region="us-west-1",
            versionId="bKufwe3zvJ3SQn3F9Z.akBkenOYl_SIz"
        )

    def test_test_event(self):
        """
        Check that the indexer does not barf when it gets an S3 test notification.
        """
        event = {
            "Records": [{
                "body": json.dumps({
                    "Message": json.dumps({
                        "Service": "Amazon S3",
                        "Event": "s3:TestEvent",
                        "Time": "2014-10-13T15:57:02.089Z",
                        "Bucket": "test-bucket",
                        "RequestId": "5582815E1AEA5ADF",
                        "HostId": "fakeGUIDhere+YstdA6Knx4Ip8EXAMPLE"
                    })
                })
            }]
        }

        index.handler(event, None)

    def test_unexpected_event(self):
        """
        Test unknown event types
        """
        # the indexer should just pass over this event without touching S3 or ES
        self._test_index_events([UNKNOWN_EVENT_TYPE], mock_elastic=False)

    def test_unknown_event_failure(self):
        """
        send unrecognizable error keys back from ES
        """
        with pytest.raises(RetryError, match="Failed to load"):
            # we don't set expected_es_calls here because the assert is never hit
            # because of the exceptions, but we do check it directly
            self._test_index_events(
                ["ObjectCreated:Put"],
                errors=True,
                status=400,
                unknown_items=True
            )
        assert self.actual_es_calls == 2, "Two failures should have called _bulk twice"

    def test_unsupported_contents(self):
        assert self._get_contents('foo.exe', '.exe') == ""
        assert self._get_contents('foo.exe.gz', '.exe.gz') == ""

    @patch(__name__ + '.index.get_content_index_bytes', return_value=100)
    def test_get_contents(self, mocked_get_content_index_bytes):
        parquet = (BASE_DIR / 'onlycolumns-c000').read_bytes()
        # mock up the responses
        size = len(parquet)
        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': size,
                'Body': BytesIO(parquet),
            }
        )
        contents = index.maybe_get_contents(
            'test-bucket',
            'some/dir/data.parquet',
            '.parquet',
            s3_client=self.s3_client,
            etag='11223344',
            size=size,
            version_id='abcde',
        )
        # test return val
        assert len(contents.encode()) == mocked_get_content_index_bytes.return_value, \
            'contents return more data than expected'
        # we know from ELASTIC_LIMIT_BYTES=1000 that column_k is the last one
        present, _, absent = ascii_lowercase.partition('l')
        for letter in present:
            col = f'column_{letter}'
            assert col in contents, f'missing column: {col}'
        for letter in absent:
            col = f'column_{letter}'
            assert col not in contents, f'missing column: {col}'

    @patch(__name__ + '.index.get_content_index_bytes', return_value=462)
    def test_get_contents_dots(self, mocked_get_content_index_bytes):
        json_ = (BASE_DIR / 'sample.json.gz').read_bytes()
        # mock up the responses
        size = len(json_)
        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': size,
                'Body': BytesIO(json_),
            }
        )
        contents = index.maybe_get_contents(
            'test-bucket',
            'some/dir.with.dots/data.another.something.json',
            '.json',
            s3_client=self.s3_client,
            etag='11223344',
            size=size,
            version_id='abcde',
            compression='gz'
        )
        dict_ = json.loads(contents)
        assert dict_["input"] == "dummy"
        assert len(dict_.keys()) == 5

    @patch.object(index, 'get_available_memory')
    def test_get_contents_large(self, get_memory_mock):
        get_memory_mock.return_value = 1
        parquet = (BASE_DIR / 'amazon-reviews-1000.snappy.parquet').read_bytes()
        # mock up the responses
        size = len(parquet)
        assert size > 1, 'supposed to test files larger than available memory'
        contents = index.maybe_get_contents(
            'test-bucket',
            'some/dir/data.parquet',
            '.parquet',
            s3_client=self.s3_client,
            etag='11223344',
            size=size,
            version_id='abcde',
        )
        # we should never touch S3 in this case and skip deserialization to avoid
        # crashing lambda by flooding memory
        assert contents == ""

    @pytest.mark.extended
    @patch('document_queue.ELASTIC_LIMIT_BYTES', 64_000)
    def test_get_contents_extended(self):
        files = (BASE_DIR / 'extended').glob('**/*-c000')
        for f in files:
            parquet = f.read_bytes()
            size = len(parquet)
            self.s3_stubber.add_response(
                method='get_object',
                service_response={
                    'Metadata': {},
                    'ContentLength': size,
                    'Body': BytesIO(parquet),
                }
            )
            contents = index.maybe_get_contents(
                'test-bucket',
                'some/dir/data.parquet',
                '.parquet',
                s3_client=self.s3_client,
                etag='11223344',
                size=size,
                version_id='abcde',
            )
            assert len(contents.encode()) <= 64_000, \
                'contents return more data than expected'

    def test_get_plain_text(self):
        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': 123,
                'Body': BytesIO(b'Hello World!\nThere is more to know.'),
            },
            expected_params={
                'Bucket': 'test-bucket',
                'Key': 'foo.txt',
                'IfMatch': 'etag',
                'Range': 'bytes=0-122',
            }
        )

        contents = index.get_plain_text(
            'test-bucket',
            'foo.txt',
            compression=None,
            etag='etag',
            version_id=None,
            s3_client=self.s3_client,
            size=123
        )
        assert contents == "Hello World!\nThere is more to know."

    def test_text_contents(self):
        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': 123,
                'Body': BytesIO(b'Hello World!'),
            },
            expected_params={
                'Bucket': 'test-bucket',
                'Key': 'foo.txt',
                'IfMatch': 'etag',
                'Range': 'bytes=0-122',
            }
        )

        assert self._get_contents('foo.txt', '.txt') == "Hello World!"

    def test_gzipped_text_contents(self):
        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': 123,
                'Body': BytesIO(compress(b'Hello World!')),
            },
            expected_params={
                'Bucket': 'test-bucket',
                'Key': 'foo.txt.gz',
                'IfMatch': 'etag',
                'Range': 'bytes=0-122',
            }
        )

        assert self._get_contents('foo.txt.gz', '.txt', compression='gz') == "Hello World!"

    def test_notebook_contents(self):
        notebook = (BASE_DIR / 'normal.ipynb').read_bytes()

        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': 123,
                'Body': BytesIO(notebook),
            },
            expected_params={
                'Bucket': 'test-bucket',
                'Key': 'foo.ipynb',
                'IfMatch': 'etag',
            }
        )

        assert "model.fit" in self._get_contents('foo.ipynb', '.ipynb')

    def test_gzipped_notebook_contents(self):
        notebook = compress((BASE_DIR / 'normal.ipynb').read_bytes())

        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': 123,
                'Body': BytesIO(notebook),
            },
            expected_params={
                'Bucket': 'test-bucket',
                'Key': 'foo.ipynb.gz',
                'IfMatch': 'etag',
            }
        )

        assert "Model results visualization" in self._get_contents('foo.ipynb.gz', '.ipynb', compression='gz')

    def test_parquet_contents(self):
        parquet = (BASE_DIR / 'amazon-reviews-1000.snappy.parquet').read_bytes()
        self.s3_stubber.add_response(
            method='get_object',
            service_response={
                'Metadata': {},
                'ContentLength': 123,
                'Body': BytesIO(parquet),
            },
            expected_params={
                'Bucket': 'test-bucket',
                'Key': 'foo.parquet',
                'IfMatch': 'etag',
            }
        )

        contents = self._get_contents('foo.parquet', '.parquet')
        size = len(contents.encode('utf-8', 'ignore'))
        assert size <= document_queue.ELASTIC_LIMIT_BYTES
        # spot check for contents
        assert "This is not even worth the money." in contents
        assert "As for results; I felt relief almost immediately." in contents
        assert "R2LO11IPLTDQDX" in contents

    # see PRE conditions in conftest.py
    @pytest.mark.extended
    def test_parquet_extended(self):
        files = (BASE_DIR / 'extended').glob('**/*.parquet')
        for f in files:
            print(f"Testing {f}")
            parquet = f.read_bytes()
            self.s3_stubber.add_response(
                method='get_object',
                service_response={
                    'Metadata': {},
                    'ContentLength': 123,
                    'Body': BytesIO(parquet),
                },
                expected_params={
                    'Bucket': 'test-bucket',
                    'Key': 'foo.parquet',
                    'IfMatch': 'etag',
                }
            )
            contents = self._get_contents('foo.parquet', '.parquet')
            size = len(contents.encode('utf-8', 'ignore'))
            assert size <= document_queue.ELASTIC_LIMIT_BYTES

    def test_get_object_tagging(self):
        bucket = "test-bucket"
        key = "test-key"
        version_id = None

        self.s3_stubber.add_response(
            method="get_object_tagging",
            service_response={
                "TagSet": [
                    {"Key": "test-key", "Value": "test-value"},
                ]
            },
            expected_params={
                "Bucket": bucket,
                "Key": key,
            },
        )

        assert index.get_object_tagging(
            s3_client=self.s3_client,
            bucket=bucket, key=key,
            version_id=version_id
        ) == {"test-key": "test-value"}

    def test_get_object_tagging_version_id(self):
        bucket = "test-bucket"
        key = "test-key"
        version_id = "test-version-id"

        self.s3_stubber.add_response(
            method="get_object_tagging",
            service_response={
                "TagSet": [
                    {"Key": "test-key", "Value": "test-value"},
                ]
            },
            expected_params={
                "Bucket": bucket,
                "Key": key,
                "VersionId": version_id,
            },
        )

        assert index.get_object_tagging(
            s3_client=self.s3_client,
            bucket=bucket, key=key,
            version_id=version_id
        ) == {"test-key": "test-value"}

    def test_get_object_tagging_access_denied(self):
        bucket = "test-bucket"
        key = "test-key"
        version_id = None

        self.s3_stubber.add_client_error(
            method="get_object_tagging",
            service_error_code="AccessDenied",
            expected_params={
                "Bucket": bucket,
                "Key": key,
            },
        )

        assert index.get_object_tagging(
            s3_client=self.s3_client,
            bucket=bucket, key=key,
            version_id=version_id
        ) is None

    def test_get_object_tagging_no_such_key(self):
        bucket = "test-bucket"
        key = "test-key"
        version_id = None

        self.s3_stubber.add_client_error(
            method="get_object_tagging",
            service_error_code="NoSuchKey",
            expected_params={
                "Bucket": bucket,
                "Key": key,
            },
        )

        with pytest.raises(botocore.exceptions.ClientError):
            assert index.get_object_tagging(
                s3_client=self.s3_client,
                bucket=bucket, key=key,
                version_id=version_id
            )


def test_extract_pptx():
    lorem = "Lorem ipsum dolor sit amet, consectetur"

    prs = pptx.Presentation()

    blank_slide_layout = prs.slide_layouts[6]
    left = top = width = height = pptx.util.Inches(1)

    slide1 = prs.slides.add_slide(blank_slide_layout)
    slide1.shapes.add_textbox(left, top, width, height).text_frame.text = lorem
    slide1.shapes.add_textbox(left, top, width, height).text_frame.text = lorem

    slide2 = prs.slides.add_slide(blank_slide_layout)
    slide2.shapes.add_textbox(left, top, width, height).text_frame.text = lorem
    slide2.shapes.add_textbox(left, top, width, height).text_frame.text = lorem

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    result = index.extract_pptx(buf, len(lorem) * 4 - 1)

    assert result == "\n".join([lorem] * 3)
